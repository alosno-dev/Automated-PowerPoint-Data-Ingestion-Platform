from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from pptx import Presentation
from io import BytesIO
import base64
from fastapi.middleware.cors import CORSMiddleware
from extractor import procesar_tabla
from etiqueta import clasificar_etiqueta
from PIL import Image


app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def emu_to_pixels(emu_value):
    """Convierte EMUs a píxeles (asumiendo 96 DPI)"""
    return int(emu_value / 9525)

@app.post("/parse-table")
async def parse_image(file: UploadFile = File(...),
    name: str = Form(...)):
    
    """
    Endpoint que recibe una imagen y devuelve la tabla OCR procesada
    """
    if not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="El archivo debe ser una imagen")

    try:
        # Leer imagen y abrir con PIL
        content = await file.read()
        img = Image.open(BytesIO(content)).convert("RGB")

        # Llamada a tu extractor.py
        tabla_resultado = await procesar_tabla(img, name)

        return tabla_resultado

    except Exception as e:
        print(f"Error procesando imagen: {e}")
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}")

@app.post("/parse-etiqueta")
async def parse_etiqueta(file: UploadFile = File(...),
    name: str = Form(...)):
    """Endpoint que recibe una imagen de etiqueta y la clasifica por color."""
    if not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="El archivo debe ser una imagen")

    try:
        # Leer imagen y abrir con PIL
        content = await file.read()
        img = Image.open(BytesIO(content)).convert("RGB")

        # Clasificar etiqueta
        etiqueta = await clasificar_etiqueta(img, name)

        return etiqueta

    except Exception as e:
        print(f"Error procesando imagen: {e}")
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}")

@app.post("/parse-pptx")
async def parse_pptx(file: UploadFile = File(...)):
    try:
        content = await file.read()
        pptx_file = BytesIO(content)
        prs = Presentation(pptx_file)

        slides_data = []

        for i, slide in enumerate(prs.slides, start=1):
            textos = []
            imagenes = []

            # Extraer textos
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    textos.append(shape.text.strip())

            # Extraer imágenes directamente en PNG con transparencia
            temp_imgs = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    try:
                        img_bytes = shape.image.blob
                        pil_img = Image.open(BytesIO(img_bytes)).convert("RGBA")

                        # Convertir directamente a base64 PNG
                        buffered = BytesIO()
                        pil_img.save(buffered, format="PNG")
                        img_base64 = base64.b64encode(buffered.getvalue()).decode("utf-8")

                        temp_imgs.append({
                            "data": img_base64,
                            "mimeType": "image/png",
                            "width": emu_to_pixels(shape.width),
                            "height": emu_to_pixels(shape.height),
                            "left": emu_to_pixels(shape.left),
                            "top": emu_to_pixels(shape.top)
                        })

                    except Exception as img_error:
                        print(f"Error extrayendo imagen en slide {i}: {img_error}")
                        continue

            # Clasificación automática de imágenes
            if temp_imgs:
                # Pegatina = imagen más pequeña
                pegatina_idx = min(range(len(temp_imgs)), key=lambda x: temp_imgs[x]["width"] * temp_imgs[x]["height"])
                temp_imgs[pegatina_idx]["categoria"] = "pegatina"

                # Tabla = imagen más a la derecha
                tabla_idx = max(range(len(temp_imgs)), key=lambda x: temp_imgs[x]["left"])
                temp_imgs[tabla_idx]["categoria"] = "tabla"

                # El resto = coche
                for idx, img in enumerate(temp_imgs):
                    if "categoria" not in img:
                        img["categoria"] = "coche"

                # Ignorar imágenes cerca de la parte inferior
                slide_height_px = emu_to_pixels(prs.slide_height)
                ignore_threshold = 100
                for img in temp_imgs:
                    if img["top"] + img["height"] >= slide_height_px - ignore_threshold:
                        img["categoria"] = "ignorar"

                # Asignar nombres finales con extensión .png
                for idx, img in enumerate(temp_imgs):
                    img["name"] = f"slide_{i}_{img['categoria']}_{idx+1}.png"
                    imagenes.append(img)

            slides_data.append({
                "slide": i,
                "textos": textos,
                "imagenes": imagenes
            })

        return {"slides": slides_data}

    except Exception as e:
        print(f"Error procesando PPTX: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}")

    finally:
        await file.close()

